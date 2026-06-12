"""Full-presentation PPTX import: slides, text, notes and theme round-trip.

The fixtures are generated with python-pptx (a real PowerPoint writer) so the
parser is exercised against genuine OOXML, not a hand-written stub.
"""

import os
import tempfile

import pytest

pptx = pytest.importorskip("pptx")
from pplib import notes, pptximport  # noqa: E402
from pplib import svgutil as S  # noqa: E402
from pptx import Presentation as PptxPresentation  # noqa: E402
from pptx.util import Emu  # noqa: E402


def _make_deck(path):
    p = PptxPresentation()
    p.slide_width = Emu(12192000)   # 16:9
    p.slide_height = Emu(6858000)

    # Slide 1: title + subtitle (title layout).
    s1 = p.slides.add_slide(p.slide_layouts[0])
    s1.shapes.title.text = "Welcome"
    s1.placeholders[1].text = "An imported deck"

    # Slide 2: title + bullet body (title+content layout).
    s2 = p.slides.add_slide(p.slide_layouts[1])
    s2.shapes.title.text = "Agenda"
    body = s2.placeholders[1].text_frame
    body.text = "First point"
    body.add_paragraph().text = "Second point"
    body.add_paragraph().text = "Third point"
    s2.notes_slide.notes_text_frame.text = "Remember to slow down here.\nBreathe."

    # Slide 3: a free-floating text box on a blank layout.
    s3 = p.slides.add_slide(p.slide_layouts[6])
    tb = s3.shapes.add_textbox(Emu(1000000), Emu(1000000), Emu(4000000), Emu(800000))
    tb.text_frame.text = "Free text box"

    p.save(path)


@pytest.fixture
def deck_path():
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    _make_deck(path)
    yield path
    os.remove(path)


_BLANK_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" '
    'xmlns:inkscape="http://www.inkscape.org/namespaces/inkscape" '
    'xmlns:sodipodi="http://sodipodi.sourceforge.net/DTD/sodipodi-0.0.dtd" '
    'width="100" height="100" viewBox="0 0 100 100">'
    '<sodipodi:namedview id="nv"/></svg>'
)


def _blank_pres():
    import inkex
    from pplib.model import Presentation

    root = inkex.load_svg(_BLANK_SVG.encode()).getroot()
    return Presentation(root)


def _slide_texts(slide):
    """All visible text strings on a slide, in document order."""
    return [S.text_content(t) for t in slide.layer.iter()
            if t.tag.endswith("}text")]


def test_import_creates_all_slides(deck_path):
    pres = _blank_pres()
    count, summary = pptximport.import_presentation(pres, deck_path)
    assert count == 3
    assert pres.slide_count() == 3
    assert "Imported 3 slides" in summary


def test_title_and_subtitle_roundtrip(deck_path):
    pres = _blank_pres()
    pptximport.import_presentation(pres, deck_path)
    texts = _slide_texts(pres.slides()[0])
    assert any("Welcome" in t for t in texts)
    assert any("An imported deck" in t for t in texts)


def test_bullets_roundtrip(deck_path):
    pres = _blank_pres()
    pptximport.import_presentation(pres, deck_path)
    joined = " ".join(_slide_texts(pres.slides()[1]))
    assert "First point" in joined
    assert "Second point" in joined
    assert "Third point" in joined


def test_notes_imported(deck_path):
    pres = _blank_pres()
    pptximport.import_presentation(pres, deck_path)
    note = notes.get_notes(pres.slides()[1])
    assert "slow down" in note
    assert "Breathe" in note


def test_free_text_box_imported(deck_path):
    pres = _blank_pres()
    pptximport.import_presentation(pres, deck_path)
    assert any("Free text box" in t for t in _slide_texts(pres.slides()[2]))


def test_append_mode_keeps_existing(deck_path, presentation):
    before = presentation.slide_count()
    count, _ = pptximport.import_presentation(
        presentation, deck_path, replace=False)
    assert presentation.slide_count() == before + count


def test_replace_mode_replaces(deck_path, presentation):
    count, _ = pptximport.import_presentation(
        presentation, deck_path, replace=True)
    assert presentation.slide_count() == count


def test_non_pptx_rejected():
    pres = _blank_pres()
    with pytest.raises(ValueError):
        pptximport.import_presentation(pres, "/tmp/whatever.odp")


def test_per_run_font_size_and_colour(tmp_path):
    """A run's explicit size and colour survive the import."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    p = PptxPresentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text = "T"
    para = s.placeholders[1].text_frame.paragraphs[0]
    para.text = "Big red"
    run = para.runs[0]
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    path = str(tmp_path / "fmt.pptx")
    p.save(path)

    pres = _blank_pres()
    pptximport.import_presentation(pres, path)
    body = next(t for t in pres.slides()[0].layer.iter()
                if t.tag.endswith("}text") and "Big red" in S.text_content(t))
    assert body.style["font-size"] == "80px"  # 40pt * 2
    assert body.style["fill"].upper() == "#C0392B"


def test_geometry_inherited_from_layout(deck_path):
    """A title with an empty spPr is positioned from the layout placeholder."""
    pres = _blank_pres()
    pptximport.import_presentation(pres, deck_path)
    title = next(t for t in pres.slides()[0].layer.iter()
                 if t.tag.endswith("}text") and "Welcome" in S.text_content(t))
    # The title-slide layout's ctrTitle sits at EMU x=685800 -> ~108px at 16:9,
    # not our default-fraction fallback (0.10 * 1920 = 192).
    assert 70 < float(title.get("x")) < 160


def test_object_z_order_preserved(tmp_path):
    """A shape drawn after a text box stays on top after import."""
    from pptx.dml.color import RGBColor

    p = PptxPresentation()
    s = p.slides.add_slide(p.slide_layouts[6])  # blank
    tb = s.shapes.add_textbox(Emu(1000000), Emu(1000000), Emu(3000000), Emu(800000))
    tb.text_frame.text = "UNDER"
    sh = s.shapes.add_shape(1, Emu(1000000), Emu(1000000), Emu(3000000), Emu(800000))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    path = str(tmp_path / "z.pptx")
    p.save(path)

    pres = _blank_pres()
    pptximport.import_presentation(pres, path)
    children = list(pres.slides()[0].layer)
    ti = next(i for i, e in enumerate(children)
              if e.tag.endswith("}text") and "UNDER" in S.text_content(e))
    ri = next(i for i, e in enumerate(children)
              if e.tag.endswith("}rect") and "000000" in (e.get("style") or ""))
    assert ti < ri  # text first (drawn earlier), shape on top -- order preserved




# A minimal but real OOXML package whose slide inherits a background *picture*
# from its slide layout -- the case branded templates use and that the importer
# must resolve through the slide -> layout chain.
import zipfile  # noqa: E402

_CT = "http://schemas.openxmlformats.org/package/2006/relationships"
_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc`\x00\x02"
    b"\x00\x00\x05\x00\x01\xe2&\x05\x9b\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _layout_bg_pptx(path):
    P, A, R = pptximport._P, pptximport._A, pptximport._R
    pres = ('<p:presentation xmlns:p="%s" xmlns:r="%s">'
            '<p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>'
            '<p:sldSz cx="12192000" cy="6858000"/></p:presentation>' % (P, R))
    pres_rels = ('<Relationships xmlns="%s">'
                 '<Relationship Id="rId1" Type="%s/slide" Target="slides/slide1.xml"/>'
                 "</Relationships>" % (_CT, _RT))
    slide = ('<p:sld xmlns:p="%s" xmlns:a="%s"><p:cSld><p:spTree>'
             '<p:sp><p:nvSpPr><p:cNvPr id="2" name="t"/><p:cNvSpPr/>'
             '<p:nvPr><p:ph type="title"/></p:nvPr></p:nvSpPr><p:spPr/>'
             '<p:txBody><a:bodyPr/><a:p><a:r><a:t>Hi</a:t></a:r></a:p></p:txBody>'
             "</p:sp></p:spTree></p:cSld></p:sld>" % (P, A))
    slide_rels = ('<Relationships xmlns="%s">'
                  '<Relationship Id="rId1" Type="%s/slideLayout" '
                  'Target="../slideLayouts/slideLayout1.xml"/>'
                  "</Relationships>" % (_CT, _RT))
    layout = ('<p:sldLayout xmlns:p="%s" xmlns:a="%s" xmlns:r="%s"><p:cSld>'
              '<p:bg><p:bgPr><a:blipFill><a:blip r:embed="rId1"/>'
              "<a:stretch><a:fillRect/></a:stretch></a:blipFill></p:bgPr></p:bg>"
              "<p:spTree/></p:cSld></p:sldLayout>" % (P, A, R))
    layout_rels = ('<Relationships xmlns="%s">'
                   '<Relationship Id="rId1" Type="%s/image" '
                   'Target="../media/image1.png"/>'
                   "</Relationships>" % (_CT, _RT))
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("ppt/presentation.xml", pres)
        zf.writestr("ppt/_rels/presentation.xml.rels", pres_rels)
        zf.writestr("ppt/slides/slide1.xml", slide)
        zf.writestr("ppt/slides/_rels/slide1.xml.rels", slide_rels)
        zf.writestr("ppt/slideLayouts/slideLayout1.xml", layout)
        zf.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", layout_rels)
        zf.writestr("ppt/media/image1.png", _PNG)


def test_layout_background_picture_resolved_per_slide(tmp_path):
    path = str(tmp_path / "branded.pptx")
    _layout_bg_pptx(path)
    pres = _blank_pres()
    count, summary = pptximport.import_presentation(pres, path)
    assert count == 1
    assert "slide backgrounds: 1" in summary
    slide = pres.slides()[0]
    # A background <image> (from the layout's p:bg picture) sits on the slide.
    imgs = [e for e in slide.layer.iter() if e.tag.endswith("}image")]
    assert imgs, "layout background picture was not imported"
    assert imgs[0].get("href", "").startswith("data:image/png;base64,")


def _inherit_size_pptx(path):
    """A deck whose body text has no run size -- it must inherit the master's."""
    P, A, R = pptximport._P, pptximport._A, pptximport._R
    pres = ('<p:presentation xmlns:p="%s" xmlns:r="%s">'
            '<p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>'
            '<p:sldSz cx="12192000" cy="6858000"/></p:presentation>' % (P, R))
    pres_rels = ('<Relationships xmlns="%s"><Relationship Id="rId1" '
                 'Type="%s/slide" Target="slides/slide1.xml"/></Relationships>'
                 % (_CT, _RT))
    # Body run carries NO sz / colour -> must inherit 28pt from the master.
    slide = ('<p:sld xmlns:p="%s" xmlns:a="%s"><p:cSld><p:spTree>'
             '<p:sp><p:nvSpPr><p:cNvPr id="2" name="b"/><p:cNvSpPr/>'
             '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/>'
             '<p:txBody><a:bodyPr/><a:p><a:r><a:t>Inherit me</a:t></a:r></a:p>'
             "</p:txBody></p:sp></p:spTree></p:cSld></p:sld>" % (P, A))
    slide_rels = ('<Relationships xmlns="%s"><Relationship Id="rId1" '
                  'Type="%s/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
                  "</Relationships>" % (_CT, _RT))
    layout = ('<p:sldLayout xmlns:p="%s" xmlns:a="%s" xmlns:r="%s"><p:cSld><p:spTree>'
              '<p:sp><p:nvSpPr><p:cNvPr id="2" name="b"/><p:cNvSpPr/>'
              '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/>'
              "<p:txBody/></p:sp></p:spTree></p:cSld></p:sldLayout>" % (P, A, R))
    layout_rels = ('<Relationships xmlns="%s"><Relationship Id="rId1" '
                   'Type="%s/slideMaster" Target="../slideMasters/slideMaster1.xml"/>'
                   "</Relationships>" % (_CT, _RT))
    master = ('<p:sldMaster xmlns:p="%s" xmlns:a="%s"><p:cSld><p:spTree>'
              '<p:sp><p:nvSpPr><p:cNvPr id="2" name="b"/><p:cNvSpPr/>'
              '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
              '<p:spPr><a:xfrm><a:off x="457200" y="1600200"/>'
              '<a:ext cx="8229600" cy="4525963"/></a:xfrm></p:spPr><p:txBody/></p:sp>'
              "</p:spTree></p:cSld><p:txStyles>"
              '<p:bodyStyle><a:lvl1pPr><a:defRPr sz="2800"/></a:lvl1pPr></p:bodyStyle>'
              "</p:txStyles></p:sldMaster>" % (P, A))
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("ppt/presentation.xml", pres)
        zf.writestr("ppt/_rels/presentation.xml.rels", pres_rels)
        zf.writestr("ppt/slides/slide1.xml", slide)
        zf.writestr("ppt/slides/_rels/slide1.xml.rels", slide_rels)
        zf.writestr("ppt/slideLayouts/slideLayout1.xml", layout)
        zf.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", layout_rels)
        zf.writestr("ppt/slideMasters/slideMaster1.xml", master)


def test_body_font_size_inherited_from_master(tmp_path):
    path = str(tmp_path / "inh.pptx")
    _inherit_size_pptx(path)
    pres = _blank_pres()
    pptximport.import_presentation(pres, path)
    body = next(t for t in pres.slides()[0].layer.iter()
                if t.tag.endswith("}text") and "Inherit me" in S.text_content(t))
    # 28pt master body style -> 56px (no run-level size on the slide).
    assert body.style["font-size"] == "56px"

